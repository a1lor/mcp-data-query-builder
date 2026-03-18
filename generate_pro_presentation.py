from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

def add_styled_slide(prs, title_text, layout_index=1):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    
    # Set background color (Dark Navy)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42) # #0F172A

    # Add a decorative top bar
    shape = slide.shapes.add_shape(
        1, # Rectangle
        0, 0, prs.slide_width, Inches(0.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(56, 189, 248) # #38BDF8 (Cyan accent)
    shape.line.visible = False

    # Style Title
    title = slide.shapes.title
    title.text = title_text
    title_text_frame = title.text_frame
    p = title_text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    return slide

def create_pro_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33) # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)

    # Main Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PROJET B : DATA QUERY BUILDER"
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.text = "MCP Server Workshop | Building Agentic Systems\nExpert: David Litvak"
    sp.font.size = Pt(24)
    sp.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
    sp.alignment = PP_ALIGN.CENTER

    # Accent Line
    line = slide.shapes.add_shape(1, Inches(4.66), Inches(5.5), Inches(4), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(56, 189, 248)
    line.line.visible = False

    # --- Slide 2: The Challenge ---
    slide = add_styled_slide(prs, "The Challenge: LLM Data Access")
    left_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(4))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Problem:"
    p.font.bold = True
    p.font.color.rgb = RGBColor(56, 189, 248)
    
    items = ["LLMs lack direct access to CSV files.", "Context windows limit data volume.", "Manual calculations are error-prone."]
    for item in items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.size = Pt(22)

    # Right Card (Visual Element)
    card = slide.shapes.add_shape(1, Inches(7), Inches(2), Inches(5), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(30, 41, 59) # Slate 800
    card.line.color.rgb = RGBColor(51, 65, 85)
    
    ctf = card.text_frame
    cp = ctf.paragraphs[0]
    cp.text = "Solution: MCP Server"
    cp.font.bold = True
    cp.alignment = PP_ALIGN.CENTER
    cp.font.color.rgb = RGBColor(56, 189, 248)
    
    cp2 = ctf.add_paragraph()
    cp2.text = "\nActs as a secure bridge between the AI model and a high-performance SQLite engine."
    cp2.font.size = Pt(18)
    cp2.font.color.rgb = RGBColor(226, 232, 240)
    cp2.alignment = PP_ALIGN.CENTER

    # --- Slide 3: Technical Architecture ---
    slide = add_styled_slide(prs, "Modern Tech Stack")
    
    cols = [
        ("CORE", "Python 3.10+\nFastMCP Framework", RGBColor(56, 189, 248)),
        ("ENGINE", "SQLite In-Memory\nStandard SQL", RGBColor(167, 139, 250)), # Purple
        ("SECURITY", "Capability Fencing\nRead-Only Mode", RGBColor(248, 113, 113)) # Red
    ]
    
    for i, (title, content, color) in enumerate(cols):
        x = Inches(1 + i * 4)
        box = slide.shapes.add_shape(1, x, Inches(2.5), Inches(3.3), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        box.line.width = Pt(2)
        box.line.color.rgb = color
        
        btf = box.text_frame
        btf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.text = title
        bp.font.bold = True
        bp.font.size = Pt(24)
        bp.font.color.rgb = color
        bp.alignment = PP_ALIGN.CENTER
        
        bp2 = btf.add_paragraph()
        bp2.text = content
        bp2.font.size = Pt(18)
        bp2.font.color.rgb = RGBColor(255, 255, 255)
        bp2.alignment = PP_ALIGN.CENTER

    # --- Slide 4: Strategic Workflow ---
    slide = add_styled_slide(prs, "Strategy: The Expert Workflow")
    
    steps = [
        ("01 RECONNAISSANCE", "Map the data terrain using schema resources."),
        ("02 ANALYSIS", "Execute targeted SQL queries and statistics."),
        ("03 SYNTHESIS", "Translate raw results into business insights.")
    ]
    
    for i, (stitle, scontent) in enumerate(steps):
        y = Inches(2 + i * 1.5)
        bar = slide.shapes.add_shape(1, Inches(1), y, Inches(11.33), Inches(1.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(30, 41, 59)
        bar.line.visible = False
        
        btf = bar.text_frame
        bp = btf.paragraphs[0]
        bp.text = stitle
        bp.font.bold = True
        bp.font.color.rgb = RGBColor(56, 189, 248)
        
        bp2 = btf.add_paragraph()
        bp2.text = scontent
        bp2.font.size = Pt(18)
        bp2.font.color.rgb = RGBColor(255, 255, 255)

    # --- Slide 5: Conclusion ---
    slide = add_styled_slide(prs, "Future-Proofing Data Agents")
    
    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Data Query Builder is just the beginning."
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    items = ["Scalable: Supports massive datasets beyond context limits.", "Safe: Integrated guardrails for enterprise environments.", "Extensible: Ready for visualization and multi-tool orchestration."]
    for item in items:
        p = tf.add_paragraph()
        p.text = f"→ {item}"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(148, 163, 184)

    prs.save("Professional_Project_B.pptx")
    print("Enhanced presentation saved as Professional_Project_B.pptx")

if __name__ == "__main__":
    create_pro_presentation()
