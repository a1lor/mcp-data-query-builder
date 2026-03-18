from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Projet B — Data Query Builder"
    subtitle.text = "MCP Server Workshop: Building Agentic Systems\nExpert: David Litvak"

    # --- Slide 2: Project Overview ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Overview"
    tf = slide.placeholders[1].text_frame
    tf.text = "Objective: Bridge the gap between LLMs and structured data."
    p = tf.add_paragraph()
    p.text = "Key features:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Load CSVs into an in-memory SQLite DB"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Enable robust SQL-based analysis"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Ensure safety via read-only constraints"
    p.level = 2

    # --- Slide 3: Technical Stack ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Technical Architecture"
    tf = slide.placeholders[1].text_frame
    tf.text = "Core Components:"
    for item in ["Python 3.10+ & FastMCP Framework", "SQLite (In-Memory) for fast querying", "Model Context Protocol (MCP) for LLM integration", "Capability Fencing for secure file access"]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 4: MCP Tools (Data Management) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "MCP Tools: Data Management"
    tf = slide.placeholders[1].text_frame
    for item in ["load_csv: Securely loads files into SQLite tables.", "list_tables: Shows current tables and row counts.", "describe_schema: Detailed view of column names and types."]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 5: MCP Tools (Analysis & Safety) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "MCP Tools: Analysis & Safety"
    tf = slide.placeholders[1].text_frame
    for item in ["run_query: Execute SELECT statements for deep analysis.", "get_statistics: Instant stats (Min, Max, Mean) on numeric columns.", "Safety First: Forbidden keyword detection (DROP, DELETE, etc.)."]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 6: Resources & Prompts ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Resources & AI Guidance"
    tf = slide.placeholders[1].text_frame
    for item in ["Resource db://schema: Direct JSON access to the database structure.", "Prompt data_query_assistant: System instructions for optimal AI-driven analysis.", "Encourages a structured 'Recon -> Analysis -> Synthesis' workflow."]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 7: Comparison: LLM vs MCP ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "LLM Alone vs. With MCP Tools"
    tf = slide.placeholders[1].text_frame
    for item in ["Accuracy: Native LLM guesses vs. SQL calculates.", "Data Volume: Prompt window limits vs. Entire CSV coverage.", "Trust: Opaque reasoning vs. Verifiable SQL queries."]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 8: Expert Workflow Strategy ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Strategy: The Expert Workflow"
    tf = slide.placeholders[1].text_frame
    for item in ["Phase 1: Reconnaissance - Mapping the data terrain.", "Phase 2: Analysis - Iterative querying and hypothesis testing.", "Phase 3: Synthesis - Transforming raw numbers into insights."]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 9: Security & Robustness ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Security & Robustness"
    tf = slide.placeholders[1].text_frame
    for item in ["Strict Read-Only Enforcement.", "Path Normalization to prevent directory traversal.", "Capability Fencing: Only DATA_ROOT subdirectories accessible."]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 10: Conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion"
    tf = slide.placeholders[1].text_frame
    tf.text = "The Data Query Builder turns an LLM into a reliable Data Analyst."
    p = tf.add_paragraph()
    p.text = "Future enhancements:"
    p.level = 1
    for item in ["Integration with external APIs (e.g., ArXiv, Weather).", "Persistent storage for cross-session analysis.", "Visualization tool support."]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 2

    prs.save("Project_B_Presentation.pptx")
    print("Presentation saved as Project_B_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
